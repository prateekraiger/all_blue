from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Configuration
OUTPUT_FILE = "ALL_BLUE_Pitch_Deck.pptx"
ASSETS_DIR = "assets"
BG_COLOR = RGBColor(15, 15, 15)  # #0f0f0f
TEXT_PRIMARY = RGBColor(232, 228, 223)  # #e8e4df
TEXT_SECONDARY = RGBColor(154, 149, 144)  # #9a9590
ACCENT_WARM = RGBColor(212, 165, 116)  # #d4a574
ACCENT_GOLD = RGBColor(201, 184, 150)  # #c9b896

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Logo
    logo_path = os.path.join(ASSETS_DIR, "logo.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(3.5), Inches(1), height=Inches(1.5))
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0), Inches(3), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ALL BLUE"
    p.font.size = Pt(60)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0), Inches(4.5), Inches(10), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "The Zenith of Curated Generosity"
    p2.font.size = Pt(24)
    p2.font.italic = True
    p2.font.color.rgb = ACCENT_WARM
    p2.alignment = PP_ALIGN.CENTER

def add_ethos_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Quote
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"We believe that every gift is a silent conversation. Our mission is to ensure that conversation is profound, precise, and permanent."'
    p.font.size = Pt(36)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

def add_friction_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    title.text_frame.text = "The Friction of Gifting"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = TEXT_PRIMARY
    
    # Points
    content = [
        ("Cognitive Overload", "In an era of infinite choice, selection becomes a chore rather than a celebration."),
        ("The Empathy Gap", "Traditional algorithms recommend based on data points, not human sentiment."),
        ("Temporal Scarcity", "The modern individual lacks the time to curate the extraordinary.")
    ]
    
    y_pos = 1.5
    for t, d in content:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(5), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = t
        p1.font.size = Pt(20)
        p1.font.color.rgb = ACCENT_GOLD
        
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_SECONDARY
        y_pos += 1.5
        
    # Image
    img_path = os.path.join(ASSETS_DIR, "hero.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6), Inches(1.5), width=Inches(3.5))

def add_solution_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(10), Inches(1))
    title.text_frame.text = "The Blue Standard"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = TEXT_PRIMARY
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Grid
    items = [
        ("Neural Concierge", "Deep-learning sentiment analysis."),
        ("Bespoke Sourcing", "Globally curated network of artisans."),
        ("Sensory Assurance", "High-fidelity AR and Voice interactions.")
    ]
    
    x_pos = 0.5
    for t, d in items:
        tb = slide.shapes.add_textbox(Inches(x_pos), Inches(3), Inches(3), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = t
        p1.font.size = Pt(24)
        p1.font.color.rgb = ACCENT_GOLD
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_SECONDARY
        p2.alignment = PP_ALIGN.CENTER
        x_pos += 3.2

def add_ai_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Image
    img_path = os.path.join(ASSETS_DIR, "ai_concept.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(4))
        
    # Title
    title = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(1))
    title.text_frame.text = "Cognitive Precision"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = TEXT_PRIMARY
    
    # Bullets
    bullets = [
        "Semantic Relationship Mapping",
        "Context-Aware Occasion Intelligence",
        "Zero-Search Discovery Workflows"
    ]
    
    tb = slide.shapes.add_textbox(Inches(5), Inches(2.5), Inches(4.5), Inches(3))
    tf = tb.text_frame
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_after = Pt(15)

def add_metrics_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    title.text_frame.text = "Measuring Excellence"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = TEXT_PRIMARY
    
    metrics = [
        ("45%", "Conversion Lift"),
        ("2.4m", "TAM (Global Luxury)"),
        ("100%", "Retention Target")
    ]
    
    x_pos = 0.5
    for val, label in metrics:
        tb = slide.shapes.add_textbox(Inches(x_pos), Inches(2.5), Inches(3), Inches(2))
        tf = tb.text_frame
        p1 = tf.paragraphs[0]
        p1.text = val
        p1.font.size = Pt(54)
        p1.font.color.rgb = ACCENT_WARM
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(16)
        p2.font.color.rgb = ACCENT_GOLD
        p2.alignment = PP_ALIGN.CENTER
        x_pos += 3.2

def add_closing_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0), Inches(2), Inches(10), Inches(1))
    title.text_frame.text = "Enter the Blue."
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = TEXT_PRIMARY
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Logo
    logo_path = os.path.join(ASSETS_DIR, "logo.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(4.25), Inches(3.5), height=Inches(1))
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(0), Inches(6), Inches(10), Inches(1))
    footer.text_frame.text = "ESTABLISHED 2024 | ALLBLUE.GIFT"
    footer.text_frame.paragraphs[0].font.size = Pt(12)
    footer.text_frame.paragraphs[0].font.color.rgb = TEXT_SECONDARY
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5) # 4:3 for classic look, or 13.33 x 7.5 for 16:9
    
    # For 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs)
    add_ethos_slide(prs)
    add_friction_slide(prs)
    add_solution_slide(prs)
    add_ai_slide(prs)
    add_metrics_slide(prs)
    add_closing_slide(prs)
    
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    main()
